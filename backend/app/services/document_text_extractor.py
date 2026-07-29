from __future__ import annotations

import logging
import shutil
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from docx import Document
from pypdf import PdfReader

logger = logging.getLogger(__name__)

try:
    from pdf2image import convert_from_path
    from PIL import Image, ImageFilter, ImageOps
    import pytesseract
except Exception:  # pragma: no cover - optional runtime dependency guard
    convert_from_path = None
    Image = None
    ImageFilter = None
    ImageOps = None
    pytesseract = None


def _ocr_image(image) -> str:
    if pytesseract is None or ImageOps is None:
        return ""

    processed = ImageOps.grayscale(image)
    processed = ImageOps.autocontrast(processed)
    if ImageFilter is not None:
        processed = processed.filter(ImageFilter.SHARPEN)
    return pytesseract.image_to_string(processed)


def _has_pdf_ocr_support() -> bool:
    return bool(convert_from_path and pytesseract and (shutil.which("pdftoppm") or shutil.which("pdftocairo")))


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1", errors="ignore")


def _clean_lines(text: str) -> str:
    lines: list[str] = []
    boilerplate = {
        "google apps",
        "main menu",
        "skip to main content",
        "keyboard shortcuts",
        "accessibility feedback",
        "shared",
        "download",
    }
    for line in (piece.strip() for piece in text.splitlines()):
        if not line:
            continue
        if line.lower() in boilerplate:
            continue
        if lines and line == lines[-1]:
            continue
        lines.append(line)
    return "\n".join(lines)


def _read_xlsx(path: Path) -> str:
    if not zipfile.is_zipfile(path):
        return ""

    namespace = {"a": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    extracted: list[str] = []

    with zipfile.ZipFile(path) as archive:
        shared_strings: list[str] = []
        if "xl/sharedStrings.xml" in archive.namelist():
            root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
            for shared in root.findall(".//a:si", namespace):
                parts = [node.text or "" for node in shared.findall(".//a:t", namespace)]
                shared_strings.append("".join(parts))

        sheet_names = sorted(
            name for name in archive.namelist() if name.startswith("xl/worksheets/sheet") and name.endswith(".xml")
        )
        for sheet_name in sheet_names:
            root = ET.fromstring(archive.read(sheet_name))
            for cell in root.findall(".//a:c", namespace):
                value = cell.findtext("a:v", default="", namespaces=namespace) or ""
                if cell.get("t") == "s" and value.isdigit():
                    index = int(value)
                    if 0 <= index < len(shared_strings):
                        value = shared_strings[index]
                elif cell.get("t") == "inlineStr":
                    value = "".join(node.text or "" for node in cell.findall(".//a:t", namespace))
                if value.strip():
                    extracted.append(value.strip())

    return "\n".join(extracted)


def _read_zip(path: Path) -> str:
    if not zipfile.is_zipfile(path):
        return ""

    extracted: list[str] = []
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(path) as archive:
            for member in archive.infolist():
                if member.is_dir():
                    continue
                suffix = Path(member.filename).suffix.lower().lstrip(".")
                if suffix not in {"pdf", "docx", "pptx", "txt", "md", "csv", "json", "log", "png", "jpg", "jpeg", "webp", "tiff", "bmp", "xlsx"}:
                    continue
                target = Path(temp_dir) / Path(member.filename).name
                target.write_bytes(archive.read(member))
                piece = extract_document_text(str(target)).strip()
                if piece:
                    extracted.append(piece)
    return "\n\n".join(extracted)


def extract_document_text(file_path: str) -> str:
    path = Path(file_path)
    if not path.exists():
        return ""

    ext = path.suffix.lower().lstrip(".")

    try:
        if ext in {"txt", "md", "csv", "json", "log"}:
            return _clean_lines(_read_text_file(path))

        if ext == "docx":
            document = Document(str(path))
            paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
            return _clean_lines("\n".join(paragraphs))

        if ext == "pptx":
            from pptx import Presentation

            presentation = Presentation(str(path))
            text_parts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            return _clean_lines("\n".join(text_parts))

        if ext == "pdf":
            reader = PdfReader(str(path))
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(page_text)

            extracted = "\n".join(text_parts).strip()
            if extracted:
                return _clean_lines(extracted)

            if not _has_pdf_ocr_support():
                return ""

            try:
                ocr_parts: list[str] = []
                for page_image in convert_from_path(str(path), dpi=220):
                    ocr_parts.append(_ocr_image(page_image))
                return _clean_lines("\n".join(part for part in ocr_parts if part.strip()))
            except Exception as exc:
                logger.debug("PDF OCR fallback skipped for %s: %s", file_path, exc)
                return ""

        if ext in {"png", "jpg", "jpeg", "webp", "tiff", "bmp"}:
            if pytesseract is None or Image is None:
                return ""
            with Image.open(str(path)) as image:
                return _clean_lines(_ocr_image(image))

        if ext == "xlsx":
            return _clean_lines(_read_xlsx(path))

        if ext == "zip":
            return _clean_lines(_read_zip(path))

    except Exception as exc:
        logger.warning("Document text extraction failed for %s: %s", file_path, exc)
        return ""

    return ""
