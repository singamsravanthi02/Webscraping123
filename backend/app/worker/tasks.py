from .celery_app import celery_app
import logging
import time
from datetime import datetime
from sqlalchemy.orm import Session
from app.db.session import SessionLocal
from app.domain.notifications.models import NotificationLog, NotificationStatus

logger = logging.getLogger(__name__)

# Basic celery task setup with retries
@celery_app.task(name="app.worker.tasks.dispatch_notification_task", bind=True, max_retries=3, default_retry_delay=30)
def dispatch_notification_task(self, log_id: int):
    """
    Worker task to actually process and send the notification via the correct channel.
    Uses robust retries (up to 3 times, with 30s delay).
    """
    db: Session = SessionLocal()
    try:
        log = db.query(NotificationLog).filter(NotificationLog.id == log_id).first()
        if not log:
            logger.error(f"Notification log {log_id} not found.")
            return

        # Mark as processing
        log.status = NotificationStatus.PROCESSING
        db.commit()

        # Simulate network latency and sending logic based on channel
        logger.info(f"Preparing to send {log.channel} to User ID {log.user_id} using template '{log.template_name}'")
        time.sleep(1) # Simulate delay
        
        # In a real app, this is where we call Twilio/SendGrid/FCM SDKs
        # For SMS: client.messages.create(to=user.phone, from_=our_phone, body=rendered_body)
        # We simulate success but allow for random network failure mocking if desired
        success = True 
        
        if success:
            log.status = NotificationStatus.SENT
            log.sent_at = datetime.utcnow()
            db.commit()
            logger.info(f"Successfully sent {log.channel} to User ID {log.user_id}")
            return True
        else:
            raise Exception("Simulated network failure")
            
    except Exception as exc:
        logger.error(f"Failed to send notification {log_id}: {str(exc)}")
        db.rollback()
        try:
            # Mark as failed in DB if we maxed out retries
            if self.request.retries >= self.max_retries:
                db_log = db.query(NotificationLog).filter(NotificationLog.id == log_id).first()
                if db_log:
                    db_log.status = NotificationStatus.FAILED
                    db_log.error_message = str(exc)
                    db.commit()
            
            # Retry
            raise self.retry(exc=exc)
        except self.MaxRetriesExceededError:
            logger.error(f"Max retries exceeded for notification {log_id}")
    finally:
        db.close()

@celery_app.task(name="app.worker.tasks.send_email")
def send_email_task(email_address: str, template: str, context: dict):
    # To be implemented
    logger.info(f"Sending {template} email to {email_address}")
    pass

@celery_app.task(name="app.worker.tasks.generate_embeddings")
def generate_embeddings_task(resource_id: str, resource_type: str):
    # To be implemented
    logger.info(f"Generating embeddings for {resource_type} {resource_id}")
    pass

@celery_app.task(name="app.worker.tasks.process_knowledge_document")
def process_knowledge_document(document_id: int, file_path: str):
    """
    Parses an uploaded document, chunks it, generates local embeddings, 
    stores in Qdrant, and flags for Question Generation.
    """
    from app.domain.knowledge.models import Document, DocumentStatus, DocumentChunk
    from app.worker.ai_pipeline import AIPipeline
    import os
    import uuid
    
    logger.info(f"Processing knowledge document ID: {document_id}")
    db: Session = SessionLocal()
    pipeline = AIPipeline()
    
    try:
        doc = db.query(Document).filter(Document.id == document_id).first()
        if not doc:
            return
            
        doc.status = DocumentStatus.PROCESSING
        db.commit()
        
        # 1. Extract Text
        ext = file_path.split('.')[-1].lower()
        raw_text = ""
        
        if ext == 'txt' or ext == 'md':
            with open(file_path, 'r', encoding='utf-8') as f:
                raw_text = f.read()
        elif ext == 'pdf':
            import pypdf
            with open(file_path, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                for page in pdf.pages:
                    raw_text += page.extract_text() + "\n"
        elif ext == 'docx':
            import docx
            d = docx.Document(file_path)
            raw_text = "\n".join([p.text for p in d.paragraphs])
        elif ext == 'pptx':
            import pptx
            p = pptx.Presentation(file_path)
            for slide in p.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        raw_text += shape.text + "\n"
                        
        # 2 & 3. Generate Embeddings & Store in Qdrant (Centralized)
        from app.domain.learning.services.qdrant_service import ingest_document, chunk_text
        
        chunks = chunk_text(raw_text)
        
        # This uses Gemini text-embedding-004 and stores in learning_materials
        ingest_document(
            document_id=doc.id,
            title=doc.title,
            text_content=raw_text,
            source_type=ext,
            subject=doc.subject,
            unit=doc.unit,
            semester=doc.semester
        )
        
        # Store metadata in PostgreSQL
        for idx, text_chunk in enumerate(chunks):
            db_chunk = DocumentChunk(
                id=f"{doc.id}{idx:04d}",
                document_id=doc.id,
                chunk_index=idx,
                char_count=len(text_chunk)
            )
            db.add(db_chunk)
            
        doc.status = DocumentStatus.COMPLETED
        db.commit()
        
        # Cleanup temp file
        if os.path.exists(file_path):
            os.remove(file_path)
            
        logger.info(f"Successfully processed document {document_id} into {len(chunks)} chunks.")
        
        # 4. Trigger Question/Content Generation
        generate_ai_content_for_document.delay(document_id)
        
    except Exception as e:
        logger.error(f"Failed to process document {document_id}: {e}")
        doc = db.query(Document).filter(Document.id == document_id).first()
        if doc:
            doc.status = DocumentStatus.FAILED
            doc.error_message = str(e)
            db.commit()
    finally:
        db.close()

@celery_app.task(name="app.worker.tasks.generate_ai_content_for_document")
def generate_ai_content_for_document(document_id: int):
    """
    Takes a processed document and generates Questions and Learning Resources via AI Agents.
    """
    from app.domain.knowledge.models import Document, GeneratedResource, LearningResourceType
    from app.domain.assessments.models import QuestionBank, QuestionType
    from app.domain.ai_orchestration.agents.content import QuestionGenerationAgent, ContentGenerationAgent
    
    logger.info(f"Generating AI content for document ID: {document_id}")
    db: Session = SessionLocal()
    
    try:
        doc = db.query(Document).filter(Document.id == document_id).first()
        if not doc or not doc.chunks:
            return
            
        from app.domain.learning.services.qdrant_service import get_document_text
        doc_text = get_document_text(document_id)
        if not doc_text:
            logger.error("No text retrieved from Qdrant for document")
            return
            
        # Truncate to first 25000 chars to avoid exceeding context window limits
        sample_text = doc_text[:25000]
        context_text = f"Document Title: {doc.title}\nSource: {doc.source}\nSubject: {doc.subject}\n\n{sample_text}"
        
        metadata = {
            "title": doc.title,
            "department": doc.department,
            "subject": doc.subject,
            "unit": doc.unit
        }
        
        # We need a system user ID for orchestration logging, usually ID 1 (Admin)
        admin_user_id = 1 
        
        q_agent = QuestionGenerationAgent()
        c_agent = ContentGenerationAgent()
        
        # 1. Generate Questions
        questions_data = q_agent.generate_questions(admin_user_id, context_text, metadata)
        for qd in questions_data:
            q_type = QuestionType.MCQ if qd.get("type") == "mcq" else QuestionType.SUBJECTIVE
            question = QuestionBank(
                topic=qd.get("topic", doc.title),
                subject=qd.get("subject", doc.subject),
                difficulty=qd.get("difficulty", 5),
                type=q_type,
                content=qd.get("content", ""),
                options=qd.get("options"),
                correct_answer=qd.get("correct_answer"),
                detailed_explanation=qd.get("detailed_explanation"),
                hints=qd.get("hints", []),
                common_mistakes=qd.get("common_mistakes", []),
                company_tags=qd.get("company_tags", []),
                bloom_level=qd.get("bloom_level"),
                placement_relevance=qd.get("placement_relevance"),
                interview_difficulty=qd.get("interview_difficulty"),
                company_difficulty=qd.get("company_difficulty"),
                estimated_time=qd.get("estimated_time"),
                marks=qd.get("marks", 1.0),
                points=qd.get("marks", 1.0)
            )
            db.add(question)
            
        # 2. Generate Resources
        resources_data = c_agent.generate_learning_resources(admin_user_id, context_text, metadata)
        
        if "flashcards" in resources_data:
            db.add(GeneratedResource(
                document_id=doc.id,
                resource_type=LearningResourceType.FLASHCARD,
                content=resources_data["flashcards"],
                topic=doc.title
            ))
            
        if "revision_notes" in resources_data:
            db.add(GeneratedResource(
                document_id=doc.id,
                resource_type=LearningResourceType.REVISION_NOTE,
                content={"markdown": resources_data["revision_notes"]},
                topic=doc.title
            ))
            
        if "cheat_sheet" in resources_data:
            db.add(GeneratedResource(
                document_id=doc.id,
                resource_type=LearningResourceType.CHEAT_SHEET,
                content={"markdown": resources_data["cheat_sheet"]},
                topic=doc.title
            ))
            
        db.commit()
        logger.info(f"Successfully generated AI content for document {document_id}")
        
    except Exception as e:
        logger.error(f"Failed to generate AI content for document {document_id}: {e}")
        db.rollback()
    finally:
        db.close()

@celery_app.task(name="app.worker.tasks.parse_resume")
def parse_resume_task(user_id: str, resume_url: str):
    from app.domain.ai_orchestration.agents.resume import ResumeAgent
    from app.domain.ai_orchestration.agents.jobs import JobMatchingAgent
    from app.domain.ai_orchestration.placement_engine import placement_engine
    from app.domain.ai_orchestration.models import StudentAIMemory
    from app.domain.users.models import User
    from app.domain.jobs.models import Job
    import os
    
    logger.info(f"Parsing resume for user {user_id}")
    db: Session = SessionLocal()
    
    try:
        user = db.query(User).filter(User.id == int(user_id)).first()
        if not user:
            return
            
        # Get local file path from url (url starts with /uploads)
        file_path = f".{resume_url}"
        if not os.path.exists(file_path):
            logger.error(f"Resume file not found: {file_path}")
            return
            
        # Extract text (assuming PDF for now)
        raw_text = ""
        ext = file_path.split('.')[-1].lower()
        if ext == 'pdf':
            import pypdf
            with open(file_path, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                for page in pdf.pages:
                    raw_text += page.extract_text() + "\n"
        elif ext == 'docx':
            import docx
            d = docx.Document(file_path)
            raw_text = "\n".join([p.text for p in d.paragraphs])
            
        if not raw_text.strip():
            logger.error("No text extracted from resume")
            return
            
        # 1. Resume Agent -> Extract Skills
        agent = ResumeAgent()
        analysis = agent.analyze_resume(user.id, raw_text)
        
        # 2. Store Skills in User
        extracted_skills = analysis.get("skills", [])
        if extracted_skills:
            # Merge with existing skills avoiding duplicates
            existing = set(user.skills or [])
            existing.update(extracted_skills)
            user.skills = list(existing)
            db.commit()
            
        # 3. Update Student AI Memory
        memory = db.query(StudentAIMemory).filter(StudentAIMemory.user_id == user.id).first()
        if not memory:
            memory = StudentAIMemory(user_id=user.id)
            db.add(memory)
            
        # Update memory context with resume details
        memory.resume_analysis = analysis
        db.commit()
        
        # 4. Recalculate Placement Readiness
        score = placement_engine.calculate_score(user.id)
        logger.info(f"Updated readiness score for user {user.id}: {score}")
        
        # 5. Trigger Job Matching Agent -> Generate Initial Recommendations
        jobs = db.query(Job).all()
        job_desc_list = [{"id": j.id, "title": j.title, "skills": j.extracted_skills} for j in jobs]
        
        if job_desc_list:
            matching_agent = JobMatchingAgent()
            profile = f"Email: {user.email}. User ID: {user.id}. Skills: {', '.join(user.skills or [])}. Goal: {user.career_goal}"
            scores = matching_agent.match_jobs(user.id, profile, job_desc_list)
            
            # Just log them for now or create notifications (the aggregate jobs task already creates notifications)
            logger.info(f"Generated {len(scores)} job recommendations for user {user.id}")
            
    except Exception as e:
        logger.error(f"Error parsing resume for user {user_id}: {e}")
        db.rollback()
    finally:
        db.close()

@celery_app.task(name="app.worker.tasks.daily_recommendation_scheduler")
def daily_recommendation_scheduler():
    """
    Runs every night via celery-beat to update readiness scores and generate adaptive plans.
    Demonstrates multi-step workflow orchestration using Celery.
    """
    from app.domain.users.models import User
    from app.domain.ai_orchestration.placement_engine import placement_engine
    from app.domain.ai_orchestration.agents.career import CareerCoachAgent
    from app.domain.ai_orchestration.models import StudentAIMemory
    
    logger.info("Starting Daily Recommendation Scheduler...")
    db: Session = SessionLocal()
    try:
        users = db.query(User).filter(User.is_active == True).all()
        coach = CareerCoachAgent()
        
        for user in users:
            # 1. Update Readiness Score
            score = placement_engine.calculate_score(user.id)
            logger.info(f"Updated readiness score for user {user.id}: {score}")
            
            # 2. Generate Adaptive Plan
            memory = db.query(StudentAIMemory).filter(StudentAIMemory.user_id == user.id).first()
            if memory:
                # Converting ORM object to dict for the prompt
                mem_dict = {
                    "weak_topics": memory.weak_topics,
                    "career_goals": memory.career_goals,
                    "score": score
                }
                # Generate and log plan silently
                coach.generate_weekly_plan(user.id, mem_dict)
                logger.info(f"Generated adaptive plan for user {user.id}")
                
        db.commit()
    except Exception as e:
        logger.error(f"Error in daily scheduler: {e}")
        db.rollback()
    finally:
        db.close()

@celery_app.task(name="app.worker.tasks.aggregate_and_process_jobs_task")
def aggregate_and_process_jobs_task():
    """
    Fetches jobs from all external aggregators, deduplicates, extracts skills via AI,
    matches them against students, and queues notifications.
    """
    from app.worker.scrapers.factory import get_all_scrapers
    from app.worker.ai_pipeline import AIPipeline
    from app.domain.jobs.models import Job, JobSource
    from app.domain.users.models import User
    from app.domain.ai_orchestration.agents.jobs import JobMatchingAgent
    from app.domain.notifications.models import NotificationLog, NotificationChannel
    from dateutil import parser
    
    logger.info("Starting Job Aggregation Pipeline...")
    db: Session = SessionLocal()
    pipeline = AIPipeline()
    matching_agent = JobMatchingAgent()
    scrapers = get_all_scrapers()
    
    new_jobs = []
    
    try:
        # 1. Fetch & Deduplicate
        for scraper in scrapers:
            try:
                raw_jobs = scraper.scrape()
            except Exception as e:
                logger.error(f"Scraper {scraper.__class__.__name__} failed: {e}. Continuing with remaining providers.")
                continue
                
            for rj in raw_jobs:
                ext_id = rj.get("external_id")
                if not ext_id:
                    continue
                
                # Check DB for duplicate
                existing = db.query(Job).filter(Job.external_id == ext_id).first()
                if not existing:
                    # Parse date safely
                    p_date = None
                    try:
                        if rj.get("posted_date"):
                            p_date = parser.parse(str(rj["posted_date"]))
                    except Exception:
                        pass
                        
                    # Extract skills via LLM
                    ai_data = pipeline.extract_job_details(rj["raw_description"])
                    
                    job = Job(
                        title=rj["title"],
                        company=rj["company"],
                        location=rj["location"],
                        salary_range=rj["salary_range"],
                        apply_link=rj["apply_link"],
                        raw_description=rj["raw_description"],
                        employment_type=rj["employment_type"],
                        posted_date=p_date,
                        external_id=ext_id,
                        source=JobSource.ARBEITNOW if "arbeitnow" in ext_id else JobSource.REMOTEOK,
                        extracted_skills=ai_data["skills"],
                        eligibility=ai_data["eligibility"],
                        ai_summary=ai_data["ai_summary"]
                    )
                    db.add(job)
                    db.flush() # get ID
                    new_jobs.append(job)
                    
        db.commit()
        logger.info(f"Aggregated {len(new_jobs)} new jobs.")
        
        # 2. Match against students and Notify
        if new_jobs:
            users = db.query(User).filter(User.is_active == True).all()
            job_desc_list = [{"id": j.id, "title": j.title, "skills": j.extracted_skills} for j in new_jobs]
            
            for user in users:
                profile = f"Email: {user.email}. User ID: {user.id} looking for tech roles."
                try:
                    scores = matching_agent.match_jobs(user.id, profile, job_desc_list)
                    for s in scores:
                        if s.get("match_score", 0) >= 80:
                            # Queue notification
                            log = NotificationLog(
                                user_id=user.id,
                                channel=NotificationChannel.IN_APP,
                                template_name="job_match_alert",
                                context_data={"job_id": s["job_id"], "score": s["match_score"]}
                            )
                            db.add(log)
                            db.flush()
                            dispatch_notification_task.delay(log.id)
                except Exception as e:
                    logger.error(f"Failed to match jobs for user {user.id}: {e}")
            db.commit()
            
    except Exception as e:
        logger.error(f"Aggregation task failed: {e}")
        db.rollback()
    finally:
        db.close()


@celery_app.task(name="app.worker.tasks.refresh_ai_job_recommendations_task")
def refresh_ai_job_recommendations_task(user_id: int | None = None):
    """
    Periodically refreshes AI job discovery queries and rankings.
    When user_id is omitted, refreshes all active student users.
    Also cleans up expired queries.
    """
    from app.domain.job_discovery.services import RecommendationEngine
    from app.domain.job_discovery.models import AIJobQuery
    from app.domain.users.models import User, UserRole
    from datetime import datetime, timezone

    db: Session = SessionLocal()
    try:
        # Cleanup expired queries
        deleted = db.query(AIJobQuery).filter(AIJobQuery.expires_at < datetime.now(timezone.utc)).delete()
        if deleted > 0:
            db.commit()
            logger.info("Deleted %s expired AI job queries", deleted)

        query = db.query(User).filter(User.is_active == True)  # noqa: E712
        if user_id:
            query = query.filter(User.id == user_id)
        else:
            query = query.filter(User.role == UserRole.STUDENT)

        users = query.all()
        refreshed = 0
        for user in users:
            try:
                engine = RecommendationEngine(db)
                engine.refresh_recommendations(user)
                refreshed += 1
            except Exception as exc:
                logger.warning("Failed to refresh AI jobs for user %s: %s", user.id, exc)
        logger.info("AI job discovery refresh completed for %s user(s)", refreshed)
        return {"refreshed": refreshed, "deleted_expired_queries": deleted}
    finally:
        db.close()
